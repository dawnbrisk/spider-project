from pptx import Presentation
from pptx.util import Inches
import os

# ----------------------- 配置区 -----------------------
folder_name = "C10"  # 自定义的变量

ppt_path = f"C:\\Users\\PCUser\\Pictures\\code\\{folder_name}-raw.pptx"
image_folder = f"C:\\Users\\PCUser\\Pictures\\code\\{folder_name}"
output_path = f"C:\\Users\\PCUser\\Pictures\\code\\{folder_name}-result.pptx"

image_pixel_width = 299
image_pixel_height = 336
dpi = 96  # 一般屏幕 DPI，除非你确认是别的

# 计算图片实际英寸
image_width_in =(image_pixel_width / dpi) * 1.27
image_height_in = (image_pixel_height / dpi) * 1.27

# 打开PPT
prs = Presentation(ppt_path)

# 获取幻灯片尺寸
slide_width = prs.slide_width.inches
slide_height = prs.slide_height.inches

# 获取图片列表（建议命名为 image01.jpg ~ image80.jpg）
image_files = sorted(
    [f for f in os.listdir(image_folder) if f.lower().endswith(('.jpg', '.jpeg', '.png'))]
)

# 校验页数
if len(image_files) != len(prs.slides):
    print(f"图片数量（{len(image_files)}）和幻灯片页数（{len(prs.slides)}）不一致！")
    exit(1)

# 插入图片（自动居中）
for i, slide in enumerate(prs.slides):
    image_path = os.path.join(image_folder, image_files[i])

    # 自动居中位置
    left = Inches((slide_width - image_width_in) / 2 + 0.3/2.54)
    top = Inches((slide_height - image_height_in) / 2 + 3 / 2.54)

    slide.shapes.add_picture(
        image_path, left, top,
        width=Inches(image_width_in),
        height=Inches(image_height_in)
    )
    print(f"插入第 {i+1} 页：{image_files[i]}")

# 保存结果
prs.save(output_path)
print(f"\n完成！新PPT已保存为：{output_path}")