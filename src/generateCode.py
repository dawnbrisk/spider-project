from pptx import Presentation
import re
import shutil
from copy import deepcopy
# 原始 PPT 路径
original_pptx = "C:\\Users\\PCUser\\Pictures\\generate_code\\example.pptx"
output_pptx = "C:\\Users\\PCUser\\Pictures\\generate_code\\result.pptx"

# 复制一份原始 PPT 作为输出
shutil.copyfile(original_pptx, output_pptx)

# 打开 PPT 文件
prs = Presentation(output_pptx)

# 获取第一页幻灯片作为模板
template_slide = prs.slides[0]

# 设置 X 和 Y 范围
x_start = 36
x_end = 37
y_start = 11
y_end = 82

# 替换文本函数
def update_text_frame_keep_style(text_frame, new_text):
    if not text_frame:
        return
    for p in text_frame.paragraphs:
        p.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = new_text

# 复制幻灯片函数
def copy_slide(prs, slide):
    slide_layout = prs.slide_layouts[6]  # 空白布局
    new_slide = prs.slides.add_slide(slide_layout)
    for shape in slide.shapes:
        new_el = deepcopy(shape.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide

# 批量生成幻灯片
for x in range(x_start, x_end + 1):
    for y in range(y_start, y_end + 1):
        location_code = f"A-{x:02d}-{y:03d}-1"

        new_slide = copy_slide(prs, template_slide)
        for shape in new_slide.shapes:
            if shape.has_text_frame:
                # 匹配类似 A-01-11-1 的原始文本
                if re.search(r"A-\d{2}-\d{3}-\d", shape.text):
                    update_text_frame_keep_style(shape.text_frame, location_code)

# 保存生成的新 PPT 文件
prs.save(output_pptx)
print(f"生成完成：共 {x_end - x_start + 1} × {y_end - y_start + 1} = {(x_end - x_start + 1)*(y_end - y_start + 1)} 张幻灯片")
print(f"保存为：{output_pptx}")
